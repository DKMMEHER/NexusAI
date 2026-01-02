import os
import sys
import logging
import base64
import io
import json
import yaml
import csv
from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
import google.generativeai as genai
from dotenv import load_dotenv

# Document Parsers
from docx import Document
from pptx import Presentation
import openpyxl
from bs4 import BeautifulSoup

load_dotenv()



# Configure Logging
logger = logging.getLogger("backend.pdf")
logger.setLevel(logging.INFO)
if not logger.handlers:
    handler = logging.StreamHandler(sys.stdout)
    handler.setFormatter(logging.Formatter("%(asctime)s %(levelname)s %(message)s"))
    logger.addHandler(handler)

# Configure Gemini
api_key = os.getenv("GEMINI_API_KEY")
if not api_key:
    logger.error("GEMINI_API_KEY not found in environment variables.")
else:
    genai.configure(api_key=api_key)

app = FastAPI(title="Universal Document Summarization Backend")

# CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

def extract_text_from_docx(file_bytes):
    doc = Document(io.BytesIO(file_bytes))
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(file_bytes):
    prs = Presentation(io.BytesIO(file_bytes))
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text_from_xlsx(file_bytes):
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
    text = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        text.append(f"--- Sheet: {sheet} ---")
        for row in ws.iter_rows(values_only=True):
            row_text = [str(cell) for cell in row if cell is not None]
            if row_text:
                text.append(", ".join(row_text))
    return "\n".join(text)

def extract_text_from_html(file_bytes):
    soup = BeautifulSoup(file_bytes, 'html.parser')
    return soup.get_text(separator='\n')

from typing import List

import uuid
from datetime import datetime

# Database Initialization
from .database import JsonDatabase, FirestoreDatabase
from auth import verify_token
from fastapi import Depends, HTTPException

project_id = os.getenv("GOOGLE_CLOUD_PROJECT")
is_cloud_run = os.getenv("K_SERVICE") is not None

if is_cloud_run and project_id:
    logger.info(f"Detected Cloud Run environment. Using Firestore (Project: {project_id}).")
    try:
        db = FirestoreDatabase(project_id)
        # Fetch initial history for RPM/RPD calc if needed, or just rely on fresh counting from DB?
        # For simplicity in this session, we might skip complex RPM calc or fetch recent on every request.
        # Let's fetch recent jobs to populate 'job_history' for RPM context if feasible, 
        # or better: query DB for recent stats.
        # ALLOWANCE: For this refactor, we will maintain a local 'job_history' cache for the session 
        # but primarily write to DB.
        # actually, get_analytics calls db.get_user_jobs.
        job_history = [] 
    except Exception as e:
        logger.error(f"Failed to initialize Firestore: {e}. Falling back to JsonDatabase.")
        db = JsonDatabase()
        job_history = db.get_user_jobs(None)
else:
    logger.info("Running locally. Using JsonDatabase.")
    db = JsonDatabase()
    job_history = db.get_user_jobs(None)

@app.post("/summarize")
async def summarize_document(
    files: List[UploadFile] = File(...), 
    prompt: str = Form(None), 
    model: str = Form("gemini-2.5-flash"),
    user_id: str = Form(None)
):
    job_id = str(uuid.uuid4())[:8]
    start_time = datetime.now()
    status = "Failed"
    response = None
    
    try:
        combined_content = []
        
        for file in files:
            filename = file.filename.lower()
            content_type = file.content_type
            file_bytes = await file.read()
            
            logger.info(f"Processing file: {filename} ({content_type})")

            # Determine processing method
            text_content = None

            # 1. Native Gemini Support (PDF)
            if filename.endswith(".pdf"):
                # For multiple PDFs, we'll convert them to parts
                combined_content.append(f"\n\n--- Document: {filename} ---\n")
                combined_content.append({
                    "mime_type": "application/pdf",
                    "data": file_bytes
                })
                continue

            # 2. Text Extraction for Office Formats
            elif filename.endswith(".docx"):
                text_content = extract_text_from_docx(file_bytes)
            elif filename.endswith(".pptx"):
                text_content = extract_text_from_pptx(file_bytes)
            elif filename.endswith(".xlsx"):
                text_content = extract_text_from_xlsx(file_bytes)
            
            # 3. Direct Text Formats
            elif filename.endswith((".txt", ".md", ".rtf", ".csv", ".json", ".xml", ".yaml", ".yml", ".log", ".py", ".js", ".html", ".htm")):
                # Try to decode as text
                try:
                    text_content = file_bytes.decode('utf-8')
                    if filename.endswith((".html", ".htm")):
                         text_content = extract_text_from_html(file_bytes)
                except UnicodeDecodeError:
                     text_content = file_bytes.decode('latin-1')

            else:
                logger.warning(f"Skipping unsupported file: {filename}")
                continue

            if text_content:
                combined_content.append(f"\n\n--- Document: {filename} ---\n")
                combined_content.append(text_content)
        
        if not combined_content:
             status = "Failed"
             return JSONResponse({"detail": "No valid content could be extracted from the uploaded files."}, status_code=400)

        # Generate Summary
        logger.info(f"Summarizing with model: {model}")
        gen_model = genai.GenerativeModel(model)
        
        user_prompt = "Summarize the following documents."
        if prompt:
            user_prompt = f"{prompt}\n\nDocuments:"
            
        # Construct the final prompt parts
        prompt_parts = [user_prompt]
        prompt_parts.extend(combined_content)
            
        response = gen_model.generate_content(prompt_parts)
        
        # Track token usage with LangSmith
        if response.usage_metadata and user_id:
            try:
                token_tracker.log_usage(
                    service="DocumentsSummarization",
                    operation="summarize_documents",
                    model=model,
                    input_tokens=response.usage_metadata.prompt_token_count,
                    output_tokens=response.usage_metadata.candidates_token_count,
                    user_id=user_id,
                    job_id=job_id or str(uuid.uuid4())
                )
            except Exception as e:
                logger.warning(f"Failed to log token usage: {e}")
        
        status = "Success"
        return JSONResponse({"summary": response.text})

    except Exception as e:
        logger.error(f"Error processing file: {e}", exc_info=True)
        print(f"ERROR: {e}", flush=True)
        status = "Failed"
        return JSONResponse({"detail": str(e)}, status_code=500)
    
    finally:
        # Calculate Metrics - approximate for now if using Firestore only, 
        # or fetch recent from DB
        # For simplicity, we'll skip complex RPM/RPD calc optimization 
        # and just save the job.
        
        # TPM (Tokens Per Minute)
        tpm = 0
        current_tokens = 0
        try:
            if hasattr(response, 'usage_metadata'):
                current_tokens = response.usage_metadata.total_token_count
        except:
            pass 
        
        # Record Job
        job_data = {
            "job_id": job_id,
            "user_id": user_id, 
            "type": "Summarization",
            "model": model,
            "rpm": 1, # Placeholder
            "tpm": current_tokens,
            "rpd": 1, # Placeholder
            "tokens": current_tokens,
            "status": status,
            "time": start_time.strftime("%Y-%m-%d %H:%M:%S")
        }
        
        # Save Job
        try:
            db.save_job(job_data)
        except Exception as e:
             logger.error(f"Failed to save job history: {e}")

@app.get("/analytics")
def get_analytics(user_id: str, token_uid: str = Depends(verify_token)):
    if token_uid != user_id:
        raise HTTPException(status_code=403, detail="User ID mismatch.")
    return db.get_user_jobs(user_id)

@app.get("/health")
def health_check_explicit():
    return {"status": "healthy"}

@app.get("/summarize")
def health_check_summarize():
    return {"status": "Documents Summarization Service Running"}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8003)
